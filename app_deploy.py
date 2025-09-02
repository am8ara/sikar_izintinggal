import streamlit as st
import os
import faiss
import json
import numpy as np
import google.generativeai as genai
from pypdf import PdfReader
from pptx import Presentation
import time
import base64

# --- Konfigurasi Halaman & API Key ---
st.set_page_config(page_title="Sistem Pakar Imigrasi", layout="wide")

# --- BACKGROUND AND STYLING ---
@st.cache_data
def get_base64_of_bin_file(bin_file):
    try:
        with open(bin_file, 'rb') as f:
            data = f.read()
        return base64.b64encode(data).decode()
    except FileNotFoundError:
        st.warning(f"File background '{bin_file}' tidak ditemukan.")
        return ""

def set_page_style(image_file):
    bin_str = get_base64_of_bin_file(image_file)
    page_bg_img_styled = f'''
    <style>
    .stApp {{
    background-image: url("data:image/jpeg;base64,{bin_str}");
    background-size: cover;
    background-repeat: no-repeat;
    background-attachment: fixed;
    }}
    [data-testid="stAppViewContainer"] > .main .block-container {{
    background-color: rgba(0, 0, 0, 0.6);
    padding: 2rem;
    border-radius: 10px;
    color: white;
    }}
    .stApp h1 {{
    color: white;
    }}
    </style>
    '''
    st.markdown(page_bg_img_styled, unsafe_allow_html=True)

# Panggil fungsi styling
set_page_style('jogja.jpeg')

# --- Judul Aplikasi ---
st.title("🇮🇩 Sistem Pakar Izin Tinggal Keimigrasian Indonesia")

# --- API KEY & FOLDER SETUP ---
try:
    GOOGLE_API_KEY = st.secrets["GOOGLE_API_KEY"]
    genai.configure(api_key=GOOGLE_API_KEY)
except (FileNotFoundError, KeyError):
    st.sidebar.warning("Google API Key tidak ditemukan. Harap masukkan di bawah.")
    GOOGLE_API_KEY = st.sidebar.text_input("Masukkan Google API Key Anda:", type="password")
    if GOOGLE_API_KEY:
        genai.configure(api_key=GOOGLE_API_KEY)
    else:
        st.info("Harap masukkan Google API Key di sidebar untuk memulai aplikasi.")
        st.stop()

FOLDER_DOKUMEN = 'dokumen_hukum'
FOLDER_QA = 'qa_databases'

# --- LOADING DATA ---
DATABASE_QA = []
if os.path.isdir(FOLDER_QA):
    for filename in os.listdir(FOLDER_QA):
        if filename.endswith('.json'):
            file_path = os.path.join(FOLDER_QA, filename)
            try:
                with open(file_path, 'r', encoding='utf-8') as f:
                    qa_pairs = json.load(f)
                    DATABASE_QA.extend(qa_pairs)
            except Exception as e:
                st.error(f"Gagal memuat file Q&A '{filename}': {e}")
else:
    st.warning(f"Folder '{FOLDER_QA}' tidak ditemukan.")

st.sidebar.success(f"Berhasil memuat {len(DATABASE_QA)} entri Q&A.")

FEW_SHOT_EXAMPLES = "--- CONTOH CARA MENJAWAB ---\nPertanyaan: Apa itu penjamin?\nJawaban: Penjamin adalah orang atau korporasi yang bertanggung jawab atas keberadaan dan kegiatan Orang Asing selama berada di Wilayah Indonesia.\n--- AKHIR CONTOH ---"

# --- FUNGSI UTAMA (DENGAN CACHING) ---
@st.cache_resource
def muat_dan_bangun_index():
    from langchain.text_splitter import RecursiveCharacterTextSplitter
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=2000, chunk_overlap=400)
    semua_potongan = []
    
    for filename in os.listdir(FOLDER_DOKUMEN):
        teks_lengkap = ""
        file_path = os.path.join(FOLDER_DOKUMEN, filename)
        try:
            if filename.endswith('.pdf'):
                reader = PdfReader(file_path)
                teks_lengkap = "".join(page.extract_text() or "" for page in reader.pages)
            elif filename.endswith('.pptx'):
                prs = Presentation(file_path)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            teks_lengkap += shape.text + "\n"
            if teks_lengkap:
                potongan_teks = text_splitter.split_text(teks_lengkap)
                for pot in potongan_teks:
                    semua_potongan.append({"sumber": filename, "konten": pot})
        except Exception as e:
            st.warning(f"Gagal memproses file {filename}: {e}")
    
    if not semua_potongan:
        st.error("Tidak ada dokumen yang bisa diproses.")
        return None, None, None, None

    konten_dokumen = [doc['konten'] for doc in semua_potongan]
    embeddings_dokumen = genai.embed_content(model="models/text-embedding-004", content=konten_dokumen, task_type="RETRIEVAL_DOCUMENT")["embedding"]
    index_dokumen = faiss.IndexFlatL2(np.array(embeddings_dokumen).shape[1])
    index_dokumen.add(np.array(embeddings_dokumen, dtype='float32'))

    texts_to_embed = []
    for item in DATABASE_QA:
        if "pertanyaan" in item and item["pertanyaan"]:
            texts_to_embed.append(item["pertanyaan"])
        elif "kata_kunci" in item and item["kata_kunci"]:
            texts_to_embed.append(item["kata_kunci"])

    index_qa = None
    if texts_to_embed:
        embeddings_qa = genai.embed_content(model="models/text-embedding-004", content=texts_to_embed, task_type="RETRIEVAL_DOCUMENT")["embedding"]
        index_qa = faiss.IndexFlatL2(np.array(embeddings_qa).shape[1])
        index_qa.add(np.array(embeddings_qa, dtype='float32'))

    return index_dokumen, semua_potongan, index_qa, DATABASE_QA

def cari_info(pertanyaan, index, bank_data, tipe, top_k=2):
    embedding_pertanyaan = np.array([genai.embed_content(model="models/text-embedding-004", content=pertanyaan, task_type="RETRIEVAL_QUERY")["embedding"]], dtype='float32')
    _, indices = index.search(embedding_pertanyaan, top_k)
    hasil = [bank_data[i] for i in indices[0]]
    if tipe == "dokumen":
        return "\n---\n".join([f"Kutipan dari {doc['sumber']}:\n{doc['konten']}" for doc in hasil])
    elif tipe == "qa":
        # Check if the keys exist before accessing them
        qa_results = []
        for doc in hasil:
            q = doc.get('pertanyaan', doc.get('kata_kunci', ''))
            a = doc.get('jawaban', doc.get('definisi', ''))
            qa_results.append(f"Pertanyaan Serupa: {q}\nJawaban yang Disarankan: {a}")
        return "\n---\n".join(qa_results)
    return ""

# --- ALUR UTAMA APLIKASI WEB ---
index_dokumen, db_dokumen, index_qa, db_qa = muat_dan_bangun_index()

if index_dokumen:
    pertanyaan_user = st.text_input("Ketik pertanyaan Anda tentang Izin Tinggal keimigrasian di sini:", "")
    if pertanyaan_user:
        with st.spinner("Menganalisis dan mencari jawaban..."):
            konteks_dokumen = cari_info(pertanyaan_user, index_dokumen, db_dokumen, "dokumen", top_k=2)
            konteks_qa = ""
            if index_qa:
                konteks_qa = cari_info(pertanyaan_user, index_qa, db_qa, "qa", top_k=1)
            
            prompt = f"""
            Anda adalah Sistem Pakar Keimigrasian Indonesia...
            (Sisa prompt Anda sama seperti sebelumnya)

            {FEW_SHOT_EXAMPLES}

            --- KONTEKS YANG DITEMUKAN ---
            [KONTEKS DARI JAWABAN SERUPA YANG SUDAH ADA]
            {konteks_qa}
            [KONTEKS DARI DOKUMEN HUKUM]
            {konteks_dokumen}
            --- AKHIR KONTEKS ---

            PERTANYAAN PENGGUNA:
            {pertanyaan_user}

            JAWABAN PAKAR:
            """
            
            model = genai.GenerativeModel('gemini-1.5-flash')
            response = model.generate_content(prompt)
            
            st.divider()
            st.subheader("Jawaban")
            st.markdown(response.text)
